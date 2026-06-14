
import streamlit as st
import requests
import pandas as pd
import numpy as np
import sqlite3
import time
import io
import os
import re
from collections import Counter

import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl

# ── Groq ─────────────────────────────────────────────────
groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY", ""))
GROQ_MODEL  = "llama-3.1-8b-instant"

# ── Page config ──────────────────────────────────────────
st.set_page_config(
    page_title="TrialIntel — Clinical Intelligence",
    page_icon="🧬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ════════════════════════════════════════════════════════
# CSS — Forces dark mode, no light mode bleed
# ════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

/* Force dark on everything */
html, body, [data-testid="stAppViewContainer"],
[data-testid="stHeader"], [data-testid="stToolbar"],
[data-testid="stDecoration"], [data-testid="stBottom"],
.main, .block-container, [class*="css"] {
  background-color: #060d1a !important;
  color: #e2e8f0 !important;
  font-family: Inter, sans-serif !important;
}

/* Remove Streamlit default white backgrounds */
[data-testid="stAppViewContainer"] > .main {
  background-color: #060d1a !important;
}
.block-container {
  padding-top: 1.2rem !important;
  padding-bottom: 2rem !important;
  max-width: 1380px !important;
  background: #060d1a !important;
}

/* Sidebar */
[data-testid="stSidebar"] {
  background: #070e1d !important;
  border-right: 1px solid #0f2040 !important;
}
[data-testid="stSidebar"] * {
  color: #94a3b8 !important;
}
[data-testid="stSidebar"] .stButton > button {
  color: white !important;
}

/* All text forced light */
p, span, div, label, h1, h2, h3, h4, h5, li {
  color: #e2e8f0 !important;
}

/* ── Hero ── */
.hero {
  background: linear-gradient(135deg,#0a1628 0%,#0d1a3a 50%,#080f20 100%);
  border: 1px solid #0f2040;
  border-radius: 16px;
  padding: 32px 36px;
  margin-bottom: 20px;
  position: relative;
  overflow: hidden;
}
.hero::before {
  content:"";
  position:absolute;
  top:-60px;right:-60px;
  width:320px;height:320px;
  background:radial-gradient(circle,rgba(59,130,246,0.06) 0%,transparent 70%);
  pointer-events:none;
}
.hero-title {
  font-size:1.9rem;
  font-weight:700;
  color:#f1f5f9 !important;
  letter-spacing:-0.5px;
  margin:0 0 6px 0;
}
.hero-sub {
  color:#475569 !important;
  font-size:0.88rem;
  margin:0;
  line-height:1.6;
}
.badge {
  display:inline-block;
  background:rgba(59,130,246,0.12);
  border:1px solid rgba(59,130,246,0.25);
  color:#60a5fa !important;
  padding:2px 10px;
  border-radius:20px;
  font-size:0.7rem;
  font-weight:500;
  margin-right:6px;
  margin-bottom:10px;
}

/* ── KPI grid ── */
.kpi-grid {
  display:grid;
  grid-template-columns:repeat(5,1fr);
  gap:12px;
  margin-bottom:20px;
}
.kpi-card {
  background:#07111f;
  border:1px solid #0f2040;
  border-radius:12px;
  padding:18px 14px;
  text-align:center;
  position:relative;
  overflow:hidden;
  transition:transform 0.2s,border-color 0.2s,box-shadow 0.2s;
}
.kpi-card:hover {
  transform:translateY(-3px);
  border-color:var(--acc,#3b82f6);
  box-shadow:0 8px 30px rgba(0,0,0,0.4);
}
.kpi-card::after {
  content:"";
  position:absolute;
  bottom:0;left:0;right:0;
  height:2px;
  background:var(--acc,#3b82f6);
}
.kpi-icon { font-size:1.2rem; margin-bottom:6px; display:block; }
.kpi-value {
  font-size:1.8rem;
  font-weight:700;
  color:var(--acc,#3b82f6) !important;
  margin:0;
  line-height:1;
  letter-spacing:-1px;
}
.kpi-label {
  font-size:0.68rem;
  color:#334155 !important;
  margin-top:5px;
  text-transform:uppercase;
  letter-spacing:0.6px;
  font-weight:600;
}
.kpi-delta {
  font-size:0.68rem;
  color:#1e3a5f !important;
  margin-top:3px;
}

/* ── Finding box ── */
.finding-box {
  background:#07111f;
  border:1px solid #0f2040;
  border-left:3px solid #3b82f6;
  border-radius:10px;
  padding:16px 20px;
  margin-bottom:18px;
  font-size:0.87rem;
  color:#94a3b8 !important;
  line-height:1.7;
}
.finding-box strong { color:#60a5fa !important; }

/* ── Shimmer banner ── */
.shimmer-banner {
  background:linear-gradient(90deg,#0d1f3c,#0f2a4a,#0d1f3c);
  background-size:200% 100%;
  animation:shimmer 1.8s infinite;
  border:1px solid #1e3a5f;
  border-radius:10px;
  padding:14px 20px;
  color:#60a5fa !important;
  font-size:0.88rem;
  font-weight:500;
  margin-bottom:16px;
  text-align:center;
}
@keyframes shimmer {
  0%{background-position:200% 0}
  100%{background-position:-200% 0}
}

/* ── Charts ── */
.chart-card {
  background:#07111f;
  border:1px solid #0f2040;
  border-radius:12px;
  padding:4px;
  margin-bottom:14px;
}

/* ── AI box ── */
.ai-box {
  background:#07111f;
  border:1px solid #0f2040;
  border-radius:14px;
  padding:28px 30px;
  line-height:1.85;
  color:#94a3b8 !important;
  font-size:0.88rem;
}
.ai-box h3 {
  color:#60a5fa !important;
  font-size:0.78rem !important;
  font-weight:600 !important;
  text-transform:uppercase;
  letter-spacing:1.2px;
  margin-top:22px !important;
  margin-bottom:8px !important;
  padding-bottom:6px;
  border-bottom:1px solid #0f2040;
}

/* ── Section label ── */
.sec-label {
  font-size:0.65rem;
  font-weight:700;
  color:#1e3a5f !important;
  text-transform:uppercase;
  letter-spacing:1.5px;
  margin:18px 0 8px 0;
  display:flex;
  align-items:center;
  gap:8px;
}
.sec-label::after {
  content:"";flex:1;
  height:1px;background:#0a1a2e;
}

/* ── Tags ── */
.tag {
  display:inline-block;
  padding:2px 9px;
  border-radius:20px;
  font-size:0.7rem;
  font-weight:500;
  margin:2px;
  border:1px solid;
  color:#60a5fa !important;
  background:rgba(59,130,246,0.08);
  border-color:rgba(59,130,246,0.2);
}

/* ── Download card ── */
.dl-card {
  background:#07111f;
  border:1px solid #0f2040;
  border-radius:12px;
  padding:22px;
}
.dl-card .dl-icon { font-size:1.8rem; margin-bottom:10px; }
.dl-card h4 { color:#e2e8f0 !important; font-size:0.92rem; margin-bottom:6px; }
.dl-card p  { color:#334155 !important; font-size:0.8rem; line-height:1.5; }

/* ── Landing card ── */
.land-card {
  background:#07111f;
  border:1px solid #0f2040;
  border-radius:12px;
  padding:24px 20px;
  text-align:center;
  transition:border-color 0.2s,transform 0.2s;
}
.land-card:hover { border-color:#1e3a5f; transform:translateY(-2px); }
.land-card .lc-icon { font-size:1.8rem; margin-bottom:12px; display:block; }
.land-card h3 { color:#cbd5e1 !important; font-size:0.88rem; font-weight:600; margin-bottom:6px; }
.land-card p  { color:#334155 !important; font-size:0.78rem; line-height:1.6; }

/* ── Tabs ── */
[data-testid="stTabs"] [data-baseweb="tab-list"] {
  background:#060d1a !important;
  border-bottom:1px solid #0f2040 !important;
  gap:2px;
}
[data-testid="stTabs"] [data-baseweb="tab"] {
  background:transparent !important;
  color:#334155 !important;
  font-size:0.83rem !important;
  font-weight:500 !important;
  border-radius:8px 8px 0 0 !important;
  padding:10px 18px !important;
}
[data-testid="stTabs"] [aria-selected="true"] {
  background:#07111f !important;
  color:#60a5fa !important;
  border-bottom:2px solid #3b82f6 !important;
}

/* ── Buttons ── */
.stButton > button {
  background:linear-gradient(135deg,#1d4ed8,#2563eb) !important;
  color:#ffffff !important;
  border:none !important;
  border-radius:10px !important;
  font-weight:600 !important;
  font-size:0.85rem !important;
  letter-spacing:0.2px !important;
  transition:all 0.2s !important;
  box-shadow:0 4px 14px rgba(37,99,235,0.35) !important;
}
.stButton > button:hover {
  transform:translateY(-1px) !important;
  box-shadow:0 6px 20px rgba(37,99,235,0.45) !important;
}
.stDownloadButton > button {
  background:linear-gradient(135deg,#065f46,#059669) !important;
  box-shadow:0 4px 14px rgba(5,150,105,0.3) !important;
  color:#ffffff !important;
}

/* ── Inputs ── */
.stTextInput input, .stNumberInput input {
  background:#07111f !important;
  border:1px solid #0f2040 !important;
  border-radius:8px !important;
  color:#e2e8f0 !important;
  font-size:0.88rem !important;
}
.stTextInput input:focus, .stNumberInput input:focus {
  border-color:#3b82f6 !important;
  box-shadow:0 0 0 2px rgba(59,130,246,0.15) !important;
}
.stSelectbox > div > div {
  background:#07111f !important;
  border:1px solid #0f2040 !important;
  border-radius:8px !important;
  color:#e2e8f0 !important;
}
.stMultiSelect > div > div {
  background:#07111f !important;
  border:1px solid #0f2040 !important;
  border-radius:8px !important;
}
[data-baseweb="select"] * { color:#e2e8f0 !important; }
[data-baseweb="menu"]  { background:#07111f !important; border:1px solid #0f2040 !important; }
[data-baseweb="option"] { background:#07111f !important; color:#e2e8f0 !important; }
[data-baseweb="option"]:hover { background:#0d1f3c !important; }

/* ── Slider ── */
.stSlider [data-baseweb="slider"] div[role="progressbar"] {
  background:linear-gradient(90deg,#1d4ed8,#3b82f6) !important;
}
[data-testid="stThumbValue"],
[data-testid="stTickBarMin"],
[data-testid="stTickBarMax"] {
  color:#475569 !important;
}

/* ── Toggle ── */
.stCheckbox label, .stToggle label { color:#94a3b8 !important; }

/* ── Progress bar ── */
.stProgress > div > div {
  background:linear-gradient(90deg,#1d4ed8,#7c3aed) !important;
  border-radius:4px !important;
}
[data-testid="stProgressBar"] {
  background:#0a1628 !important;
}

/* ── Dataframe ── */
[data-testid="stDataFrame"] {
  border:1px solid #0f2040 !important;
  border-radius:10px !important;
  overflow:hidden !important;
}
.dvn-scroller { background:#07111f !important; }
.col-header-cell-comp, .cell-comp {
  background:#07111f !important;
  color:#94a3b8 !important;
  border-color:#0f2040 !important;
}

/* ── Scrollbar ── */
::-webkit-scrollbar { width:5px; height:5px; }
::-webkit-scrollbar-track { background:#060d1a; }
::-webkit-scrollbar-thumb { background:#0f2040; border-radius:3px; }
::-webkit-scrollbar-thumb:hover { background:#1e3a5f; }

/* ── Caption / small text ── */
.stCaption, small { color:#334155 !important; }

/* ── Hide Streamlit branding ── */
#MainMenu, footer, header { visibility:hidden !important; }
[data-testid="stToolbar"] { display:none !important; }
</style>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════
# CONSTANTS
# ════════════════════════════════════════════════════════

CATEGORY_CONDITIONS = {
    "Custom":            [],
    "💊 Pharma":         ["Diabetes","Hypertension","Depression","Asthma","COPD"],
    "🔬 Biotech":        ["Cancer Immunotherapy","Gene Therapy","Cell Therapy","mRNA"],
    "🩺 MedTech":        ["Cardiac Stent","Hip Replacement","Pacemaker","Cochlear Implant"],
    "🧬 Biologics":      ["Rheumatoid Arthritis","Multiple Sclerosis","Psoriasis","Crohn"],
    "🌿 Nutraceuticals": ["Vitamin D","Omega-3","Probiotic","Melatonin","Magnesium"],
}

COLORS = {
    "Industry":"#3B82F6","Government / NIH":"#F59E0B",
    "Government / Federal":"#D97706","Government / Other":"#B45309",
    "Academic / Network":"#10B981","Academic / Other":"#059669",
    "Individual":"#EF4444",
}
STATUS_COLORS = {
    "Completed":"#10B981","Recruiting":"#3B82F6",
    "Active, Not Recruiting":"#6366F1","Terminated":"#EF4444",
    "Withdrawn":"#F87171","Suspended":"#F59E0B",
    "Not Yet Recruiting":"#60A5FA","Unknown":"#9CA3AF",
}
LAYOUT = dict(
    plot_bgcolor="#07111f",
    paper_bgcolor="#07111f",
    font=dict(color="#475569", family="Inter, sans-serif", size=11),
    title_font=dict(size=13, color="#cbd5e1", family="Inter, sans-serif"),
)


# ════════════════════════════════════════════════════════
# PIPELINE
# ════════════════════════════════════════════════════════

def fetch_trials(condition, max_results=500):
    BASE_URL = "https://clinicaltrials.gov/api/v2/studies"
    FIELDS = ",".join([
        "NCTId","BriefTitle","Phase","OverallStatus",
        "StartDate","PrimaryCompletionDate","EnrollmentCount",
        "LeadSponsorName","LeadSponsorClass","LocationCountry",
        "StudyType","WhyStopped","Condition",
        "InterventionType","InterventionName",
    ])
    all_rows, next_token, page = [], None, 1
    prog = st.progress(0, text="")
    info = st.empty()

    while len(all_rows) < max_results:
        params = {
            "query.cond": condition, "fields": FIELDS,
            "pageSize": 100, "format": "json", "countTotal": "true",
        }
        if next_token:
            params["pageToken"] = next_token
        try:
            r = requests.get(BASE_URL, params=params, timeout=30)
            r.raise_for_status()
        except Exception as e:
            st.error(f"API error: {e}"); break

        data = r.json()
        if page == 1:
            total = data.get("totalCount", max_results)
            info.caption(f"Found {total:,} trials — fetching up to {max_results:,}")

        studies = data.get("studies", [])
        if not studies: break

        for study in studies:
            all_rows.append(_parse_study(study))

        pct = min(len(all_rows) / max_results, 1.0)
        prog.progress(pct, text=f"Fetching... {len(all_rows):,} trials")
        next_token = data.get("nextPageToken")
        if not next_token: break
        page += 1
        time.sleep(0.2)

    prog.empty(); info.empty()
    return pd.DataFrame(all_rows)


def _parse_study(study):
    proto       = study.get("protocolSection", {})
    id_mod      = proto.get("identificationModule", {})
    status_mod  = proto.get("statusModule", {})
    design_mod  = proto.get("designModule", {})
    sponsor_mod = proto.get("sponsorCollaboratorsModule", {})
    cond_mod    = proto.get("conditionsModule", {})
    interv_mod  = proto.get("armsInterventionsModule", {})
    loc_mod     = proto.get("contactsLocationsModule", {})

    phases    = design_mod.get("phases", [])
    phase     = phases[0] if phases else "NOT_APPLICABLE"
    phase     = phase.replace("PHASE","Phase ").replace("_"," ").title()
    locations = loc_mod.get("locations", [])
    countries = list({l.get("country","") for l in locations if l.get("country")})
    intervs   = interv_mod.get("interventions", [])
    i_types   = list({i.get("type","") for i in intervs if i.get("type")})

    return {
        "nct_id":          id_mod.get("nctId",""),
        "title":           id_mod.get("briefTitle",""),
        "phase":           phase,
        "status":          status_mod.get("overallStatus","UNKNOWN").replace("_"," ").title(),
        "study_type":      design_mod.get("studyType",""),
        "start_date":      status_mod.get("startDateStruct",{}).get("date",""),
        "completion_date": status_mod.get("primaryCompletionDateStruct",{}).get("date",""),
        "enrollment":      design_mod.get("enrollmentInfo",{}).get("count",0),
        "sponsor":         sponsor_mod.get("leadSponsor",{}).get("name",""),
        "sponsor_type":    sponsor_mod.get("leadSponsor",{}).get("class",""),
        "countries":       " | ".join(countries) if countries else "Not specified",
        "country_count":   len(countries),
        "why_stopped":     status_mod.get("whyStopped",""),
        "conditions":      " | ".join(cond_mod.get("conditions",[])[:3]),
        "interv_types":    " | ".join(i_types),
    }


def clean_trials(df):
    df = df.copy()
    for col in ["phase","status","sponsor","sponsor_type",
                "study_type","countries","why_stopped","interv_types"]:
        if col in df.columns:
            df[col] = df[col].astype(str).replace("nan","").replace("None","")

    df["start_date"]      = pd.to_datetime(df["start_date"],      errors="coerce")
    df["completion_date"] = pd.to_datetime(df["completion_date"], errors="coerce")
    df["enrollment"]      = pd.to_numeric(df["enrollment"],       errors="coerce").fillna(0).astype(int)
    df["start_year"]      = df["start_date"].dt.year.astype("Int64")

    diff = (df["completion_date"] - df["start_date"]).dt.days
    df["duration_months"] = (diff / 30.44).round(1)
    df.loc[df["duration_months"] < 0,   "duration_months"] = np.nan
    df.loc[df["duration_months"] > 300, "duration_months"] = np.nan

    df["status_clean"]  = df["status"].str.strip().str.title()
    df["is_terminated"] = (df["status_clean"] == "Terminated").astype(int)
    df["primary_country"] = df["countries"].str.split(" | ").str[0].str.strip()
    df["primary_country"] = df["primary_country"].replace(
        {"Not specified":np.nan,"":np.nan,"nan":np.nan}
    )
    smap = {
        "INDUSTRY":"Industry","NIH":"Government / NIH",
        "FED":"Government / Federal","OTHER_GOV":"Government / Other",
        "NETWORK":"Academic / Network","INDIV":"Individual","OTHER":"Academic / Other",
    }
    df["sponsor_category"] = (df["sponsor_type"].str.upper().str.strip()
                               .map(smap).fillna("Academic / Other"))
    bins   = [-1,0,49,199,999,4999,float("inf")]
    labels = ["Unknown","Very small (<50)","Small (50-199)",
              "Medium (200-999)","Large (1k-4.9k)","Very large (5k+)"]
    df["enrollment_bucket"] = pd.cut(df["enrollment"], bins=bins, labels=labels)
    df["is_industry"]       = (df["sponsor_category"] == "Industry").astype(int)
    df = df.dropna(subset=["nct_id"]).drop_duplicates(subset=["nct_id"], keep="first")
    df = df.sort_values("start_date", ascending=False).reset_index(drop=True)
    return df


def apply_filters(df, phases, statuses, sponsor_cats, year_range, min_enroll):
    d = df.copy()
    if phases:       d = d[d["phase"].isin(phases)]
    if statuses:     d = d[d["status_clean"].isin(statuses)]
    if sponsor_cats: d = d[d["sponsor_category"].isin(sponsor_cats)]
    if year_range:
        d = d[(d["start_year"]>=year_range[0])&(d["start_year"]<=year_range[1])]
    if min_enroll > 0: d = d[d["enrollment"] >= min_enroll]
    return d


def compute_summary(df, condition):
    conn = sqlite3.connect(":memory:")
    df.to_sql("t", conn, index=False, if_exists="replace")
    def sql(q): return pd.read_sql_query(q, conn)

    r = sql("""SELECT COUNT(*) as total, SUM(is_terminated) as terminated,
                      ROUND(SUM(is_terminated)*100.0/COUNT(*),1) as dropout_rate,
                      ROUND(AVG(CASE WHEN enrollment>0 THEN enrollment END),0) as avg_enrollment,
                      COUNT(DISTINCT sponsor) as unique_sponsors,
                      COUNT(DISTINCT primary_country) as unique_countries,
                      MIN(start_year) as earliest_year, MAX(start_year) as latest_year
               FROM t""").iloc[0]

    top_sp = sql("""SELECT sponsor, COUNT(*) as n FROM t
                    WHERE sponsor!="" AND sponsor!="Unknown"
                    GROUP BY sponsor ORDER BY n DESC LIMIT 1""").iloc[0]
    top_ph = sql("""SELECT phase, COUNT(*) as n FROM t
                    GROUP BY phase ORDER BY n DESC LIMIT 1""").iloc[0]
    top_ct = sql("""SELECT primary_country, COUNT(*) as n FROM t
                    WHERE primary_country IS NOT NULL AND primary_country!="nan"
                    GROUP BY primary_country ORDER BY n DESC LIMIT 1""")
    top_r  = sql("""SELECT why_stopped, COUNT(*) as n FROM t
                    WHERE is_terminated=1 AND why_stopped!="" AND why_stopped!="nan"
                    GROUP BY why_stopped ORDER BY n DESC LIMIT 1""")
    sp     = sql("""SELECT ROUND(SUM(is_industry)*100.0/COUNT(*),1) as ind,
                           ROUND(SUM(1-is_industry)*100.0/COUNT(*),1) as acad
                    FROM t""").iloc[0]
    pk     = sql("""SELECT start_year, COUNT(*) as n FROM t
                    WHERE start_year IS NOT NULL
                    GROUP BY start_year ORDER BY n DESC LIMIT 1""").iloc[0]

    return {
        "condition":          condition,
        "total_trials":       int(r["total"]),
        "terminated":         int(r["terminated"]),
        "dropout_rate":       float(r["dropout_rate"]),
        "avg_enrollment":     int(r["avg_enrollment"]) if pd.notna(r["avg_enrollment"]) else 0,
        "unique_sponsors":    int(r["unique_sponsors"]),
        "unique_countries":   int(r["unique_countries"]),
        "year_range":         f"{int(r['earliest_year'])}–{int(r['latest_year'])}",
        "top_sponsor":        str(top_sp["sponsor"]),
        "top_sponsor_trials": int(top_sp["n"]),
        "top_phase":          str(top_ph["phase"]),
        "top_country":        str(top_ct.iloc[0]["primary_country"]) if len(top_ct)>0 else "N/A",
        "top_stop_reason":    str(top_r.iloc[0]["why_stopped"])[:80] if len(top_r)>0 else "N/A",
        "industry_pct":       float(sp["ind"]),
        "academic_pct":       float(sp["acad"]),
        "peak_year":          int(pk["start_year"]),
        "peak_year_trials":   int(pk["n"]),
    }


def get_ai_analysis(s):
    prompt = f"""You are a senior healthcare consultant.
Analyze this clinical trial data and write a consulting-grade report.

CONDITION: {s["condition"].upper()}
Trials: {s["total_trials"]:,} | Range: {s["year_range"]} | Dropout: {s["dropout_rate"]}%
Avg enrollment: {s["avg_enrollment"]:,} | Sponsors: {s["unique_sponsors"]:,} | Countries: {s["unique_countries"]}
Top phase: {s["top_phase"]} | Top sponsor: {s["top_sponsor"]} ({s["top_sponsor_trials"]} trials)
Industry: {s["industry_pct"]}% | Peak year: {s["peak_year"]} ({s["peak_year_trials"]} trials)
Top termination reason: {s["top_stop_reason"]}

Write exactly 4 sections with these headers:

### Market Maturity Assessment
### Competitive Dynamics
### Risk Profile
### Strategic Recommendations

Each section 3-4 sentences. Data-driven. No filler."""

    r = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[{"role":"user","content":prompt}],
        max_tokens=1200, temperature=0.7,
    )
    return r.choices[0].message.content


# ════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════

def chart_phase(df):
    d = df["phase"].value_counts().reset_index()
    d.columns = ["phase","trials"]
    fig = px.bar(d, x="phase", y="trials", color="trials",
                 color_continuous_scale=[[0,"#0f2040"],[1,"#3b82f6"]],
                 text="trials", title="Phase Distribution",
                 labels={"phase":"","trials":"Trials"})
    fig.update_traces(textposition="outside",
                      textfont=dict(color="#475569",size=10),
                      marker_line_width=0)
    fig.update_layout(**LAYOUT, height=300,
        margin=dict(t=40,b=25,l=30,r=20),
        xaxis=dict(categoryorder="total descending",showgrid=False,
                   tickfont=dict(size=10),color="#334155"),
        yaxis=dict(showgrid=True,gridcolor="#0a1628",color="#334155"),
        coloraxis_showscale=False)
    return fig

def chart_sponsors(df):
    d = (df.groupby(["sponsor","sponsor_category"])
           .size().reset_index(name="trials")
           .sort_values("trials",ascending=False).head(12))
    d["s"] = d["sponsor"].str[:34]
    fig = px.bar(d, x="trials", y="s",
                 color="sponsor_category", color_discrete_map=COLORS,
                 orientation="h", text="trials", title="Top Sponsors",
                 labels={"trials":"","s":"","sponsor_category":""})
    fig.update_traces(textposition="outside",
                      textfont=dict(size=10),marker_line_width=0)
    fig.update_layout(**LAYOUT, height=400,
        margin=dict(t=40,b=25,l=210,r=60),
        yaxis=dict(categoryorder="total ascending",showgrid=False,
                   tickfont=dict(size=10),color="#475569"),
        xaxis=dict(showgrid=True,gridcolor="#0a1628",color="#334155"),
        legend=dict(orientation="h",yanchor="bottom",y=1.01,
                    bgcolor="rgba(0,0,0,0)",font=dict(size=9),
                    title_text=""))
    return fig

def chart_status(df):
    d = df["status_clean"].value_counts().reset_index()
    d.columns = ["status","trials"]
    fig = px.pie(d, names="status", values="trials", hole=0.55,
                 color="status", color_discrete_map=STATUS_COLORS,
                 title="Trial Status")
    fig.update_traces(textposition="outside",textinfo="percent+label",
                      textfont_size=10,pull=[0.015]*len(d),
                      marker=dict(line=dict(color="#060d1a",width=2)))
    fig.update_layout(**LAYOUT, height=320,
        margin=dict(t=40,b=45,l=45,r=45), showlegend=False)
    return fig

def chart_map(df):
    all_c = []
    for row in df["countries"].dropna():
        for c in str(row).split(" | "):
            c = c.strip()
            if c and c not in ("","nan","Not specified"):
                all_c.append(c)
    cdf = pd.DataFrame(Counter(all_c).most_common(60),columns=["country","trials"])
    fig = px.choropleth(cdf, locations="country", locationmode="country names",
                        color="trials",
                        color_continuous_scale=[[0,"#0a1628"],[0.5,"#1e3a5f"],[1,"#3b82f6"]],
                        title="Geographic Distribution",labels={"trials":"Trials"})
    fig.update_layout(**LAYOUT, height=330,
        margin=dict(t=40,b=10,l=10,r=10),
        geo=dict(showframe=False,bgcolor="#07111f",
                 landcolor="#0d1f3c",showocean=True,
                 oceancolor="#060d1a",showcoastlines=True,
                 coastlinecolor="#0f2040",
                 projection_type="natural earth"),
        coloraxis_colorbar=dict(
            thickness=8,len=0.5,
            tickfont=dict(color="#334155",size=9),
            title=dict(text="",font=dict(color="#334155"))))
    return fig

def chart_timeline(df):
    y = (df[df["start_year"].notna()]
           .groupby("start_year")
           .agg(total=("nct_id","count"),industry=("is_industry","sum"))
           .reset_index())
    y["academic"] = y["total"] - y["industry"]
    y = y[(y["start_year"]>=2000)&(y["start_year"]<=2024)]
    y["start_year"] = y["start_year"].astype(int)
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=y["start_year"],y=y["industry"],name="Industry",
        fill="tozeroy",stackgroup="one",
        line=dict(color="#3b82f6",width=1.5),
        fillcolor="rgba(59,130,246,0.35)"))
    fig.add_trace(go.Scatter(
        x=y["start_year"],y=y["academic"],name="Academic/Gov",
        fill="tonexty",stackgroup="one",
        line=dict(color="#10b981",width=1.5),
        fillcolor="rgba(16,185,129,0.25)"))
    fig.update_layout(**LAYOUT, height=300,
        margin=dict(t=40,b=30,l=50,r=20),
        title="Activity Over Time",
        xaxis=dict(showgrid=False,dtick=4,color="#334155",tickfont=dict(size=10)),
        yaxis=dict(showgrid=True,gridcolor="#0a1628",color="#334155",title="New Trials"),
        hovermode="x unified",
        legend=dict(orientation="h",yanchor="bottom",y=1.01,
                    bgcolor="rgba(0,0,0,0)",font=dict(size=10),title_text=""))
    return fig

def chart_dropout(df):
    d = (df.groupby("phase")
           .agg(trials=("nct_id","count"),terminated=("is_terminated","sum"))
           .reset_index())
    d["dp"] = (d["terminated"]/d["trials"]*100).round(1)
    d = d[d["trials"]>5].sort_values("dp",ascending=False)
    fig = make_subplots(specs=[[{"secondary_y":True}]])
    fig.add_trace(go.Bar(x=d["phase"],y=d["trials"],name="Trials",
                         marker_color="#0f2040",
                         marker_line_color="#1e3a5f",marker_line_width=1),
                  secondary_y=False)
    fig.add_trace(go.Scatter(x=d["phase"],y=d["dp"],name="Dropout %",
                             line=dict(color="#ef4444",width=2),
                             mode="lines+markers",
                             marker=dict(size=6,color="#ef4444",
                                         line=dict(color="#060d1a",width=2))),
                  secondary_y=True)
    fig.update_layout(**LAYOUT, height=300,
        margin=dict(t=40,b=30,l=50,r=55),
        title="Volume vs Dropout Rate")
    fig.update_yaxes(showgrid=True,gridcolor="#0a1628",color="#334155",secondary_y=False)
    fig.update_yaxes(showgrid=False,ticksuffix="%",color="#ef4444",secondary_y=True)
    return fig

def chart_enrollment(df):
    d = df[df["enrollment"]>0].copy()
    if len(d) == 0: return go.Figure()
    fig = px.box(d, x="phase", y="enrollment",
                 color="phase",
                 color_discrete_sequence=["#3b82f6","#10b981","#f59e0b",
                                           "#8b5cf6","#ef4444","#06b6d4"],
                 title="Enrollment Distribution by Phase",
                 labels={"phase":"","enrollment":"Participants"})
    fig.update_layout(**LAYOUT, height=300,
        margin=dict(t=40,b=25,l=55,r=20),
        xaxis=dict(showgrid=False,color="#334155",tickfont=dict(size=10)),
        yaxis=dict(showgrid=True,gridcolor="#0a1628",color="#334155"),
        showlegend=False)
    return fig


# ════════════════════════════════════════════════════════
# EXPORTS
# ════════════════════════════════════════════════════════

def build_excel(df, s):
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, sheet_name="All Trials", index=False)
        phase_df = df.groupby("phase").agg(
            trials=("nct_id","count"),
            avg_enrollment=("enrollment","mean"),
            terminated=("is_terminated","sum")).reset_index()
        phase_df["dropout_pct"] = (phase_df["terminated"]/phase_df["trials"]*100).round(1)
        phase_df.to_excel(w, sheet_name="By Phase", index=False)
        (df.groupby("sponsor").agg(trials=("nct_id","count"))
           .sort_values("trials",ascending=False).head(50)
           .to_excel(w, sheet_name="Top Sponsors"))
        (df.groupby("primary_country").agg(trials=("nct_id","count"))
           .sort_values("trials",ascending=False).head(30)
           .to_excel(w, sheet_name="Geography"))
    buf.seek(0); return buf


def build_pptx(df, s):
    DARK=RGBColor(6,13,26);   CARD=RGBColor(7,17,31)
    BDR=RGBColor(15,32,64);   BLUE=RGBColor(59,130,246)
    BLDK=RGBColor(29,78,216); GREEN=RGBColor(16,185,129)
    RED=RGBColor(239,68,68);  AMBER=RGBColor(245,158,11)
    PURP=RGBColor(139,92,246);WHITE=RGBColor(255,255,255)
    GRAY=RGBColor(51,65,85);  LIGHT=RGBColor(148,163,184)

    prs=Presentation()
    prs.slide_width=Inches(13.33)
    prs.slide_height=Inches(7.5)

    def ns():
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        bg=sl.background.fill;bg.solid();bg.fore_color.rgb=DARK
        return sl
    def rc(sl,x,y,w,h,f=None):
        sh=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
        sh.fill.solid() if f else sh.fill.background()
        if f: sh.fill.fore_color.rgb=f
        sh.line.fill.background();return sh
    def tx(sl,t,x,y,w,h,sz=13,bold=False,col=WHITE,
           al=PP_ALIGN.LEFT,it=False):
        tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
        tf=tb.text_frame;tf.word_wrap=True
        p=tf.paragraphs[0];p.alignment=al
        r=p.add_run();r.text=str(t)
        r.font.size=Pt(sz);r.font.bold=bold
        r.font.italic=it;r.font.color.rgb=col
    def kc(sl,x,y,v,lb,vc=BLUE):
        rc(sl,x,y,2.9,1.4,CARD)
        tx(sl,v,x+0.1,y+0.1,2.7,0.8,sz=34,bold=True,col=vc,al=PP_ALIGN.CENTER)
        tx(sl,lb,x+0.1,y+0.9,2.7,0.4,sz=11,col=GRAY,al=PP_ALIGN.CENTER)
    def sf(v,n=45): v=str(v); return v[:n]+"..." if len(v)>n else v

    s1=ns()
    rc(s1,0,0,13.33,0.07,BLDK)
    tx(s1,f"Clinical Trial Landscape: {s['condition'].title()}",0.5,0.18,12,0.82,sz=33,bold=True)
    tx(s1,f"Competitive Intelligence  ·  {s['total_trials']:,} trials  ·  {s['year_range']}  ·  ClinicalTrials.gov",
       0.5,1.05,12.3,0.38,sz=12,it=True,col=GRAY)
    rc(s1,0.5,1.48,12.33,0.02,BDR)
    kc(s1,0.40,1.65,f"{s['total_trials']:,}","Trials",BLUE)
    kc(s1,3.52,1.65,f"{s['dropout_rate']}%","Dropout",RED)
    kc(s1,6.64,1.65,f"{s['unique_sponsors']:,}","Sponsors",GREEN)
    kc(s1,9.76,1.65,f"{s['unique_countries']}","Countries",AMBER)
    rc(s1,0.5,3.25,12.33,1.4,RGBColor(7,17,31))
    tx(s1,"KEY FINDING",0.75,3.35,4,0.3,sz=10,bold=True,col=BLUE)
    tx(s1,(f"{s['top_phase']} trials dominate. {sf(s['top_sponsor'],40)} leads "
           f"with {s['top_sponsor_trials']} trials. "
           f"Industry: {s['industry_pct']}%. Peak: {s['peak_year']}."),
       0.75,3.72,12.0,0.88,sz=13,col=LIGHT)
    for i,(st2,col) in enumerate([
        (f"Avg enrollment: {s['avg_enrollment']:,}",BLUE),
        (f"Top country: {s['top_country']}",AMBER),
        (f"Peak year: {s['peak_year']}",GREEN),
    ]):
        tx(s1,st2,0.5+i*4.3,4.82,4.1,0.42,sz=12,bold=True,col=col)
    tx(s1,"Confidential  ·  ClinicalTrials.gov  ·  1 of 3",
       0.5,7.1,12.33,0.28,sz=9,it=True,col=GRAY,al=PP_ALIGN.RIGHT)

    s2=ns()
    rc(s2,0,0,13.33,0.07,BLUE)
    tx(s2,"Competitive Sponsor Landscape",0.5,0.18,12,0.7,sz=27,bold=True)
    tx(s2,"Top 10 sponsors · trial volume · dropout % color-coded",
       0.5,0.9,12,0.32,sz=11,it=True,col=GRAY)
    rc(s2,0.5,1.24,12.33,0.02,BDR)
    for hdr,x,w,al in [
        ("SPONSOR",0.55,5.5,PP_ALIGN.LEFT),("TRIALS",6.2,1.4,PP_ALIGN.CENTER),
        ("DROPOUT",7.7,1.4,PP_ALIGN.CENTER),("AVG ENROLL",9.2,1.6,PP_ALIGN.CENTER),
        ("TYPE",10.9,2.1,PP_ALIGN.CENTER),
    ]:
        tx(s2,hdr,x,1.3,w,0.3,sz=9,bold=True,col=GRAY,al=al)
    c2=sqlite3.connect(":memory:")
    df.to_sql("t",c2,index=False,if_exists="replace")
    t10=__import__("pandas").read_sql_query("""
        SELECT sponsor,sponsor_category,COUNT(*) as trials,
               ROUND(SUM(is_terminated)*100.0/COUNT(*),1) as dp,
               ROUND(AVG(CASE WHEN enrollment>0 THEN enrollment END),0) as ae
        FROM t WHERE sponsor!="" AND sponsor!="Unknown"
        GROUP BY sponsor,sponsor_category ORDER BY trials DESC LIMIT 10
    """,c2)
    cc_map={"Industry":BLUE,"Government / NIH":AMBER,"Government / Federal":AMBER,
            "Government / Other":AMBER,"Academic / Network":GREEN,
            "Academic / Other":GREEN,"Individual":RED}
    mx=int(t10["trials"].max())
    for i,row in t10.iterrows():
        y=1.65+i*0.52
        if i%2==0: rc(s2,0.5,y-0.04,12.33,0.5,CARD)
        cc=cc_map.get(str(row["sponsor_category"]),PURP)
        bw=max(0.05,int(row["trials"])/mx*5.5)
        rc(s2,0.55,y+0.06,bw,0.26,
           RGBColor(min(255,int(cc[0]*0.2)),min(255,int(cc[1]*0.2)),min(255,int(cc[2]*0.2))))
        tx(s2,sf(row["sponsor"],48),0.65,y+0.07,5.4,0.26,sz=10,bold=(i==0))
        dp=row["dp"]
        dc=RED if dp>20 else (AMBER if dp>10 else GREEN)
        tx(s2,str(int(row["trials"])),6.2,y+0.07,1.4,0.26,sz=11,bold=True,col=cc,al=PP_ALIGN.CENTER)
        tx(s2,f"{dp}%",7.7,y+0.07,1.4,0.26,sz=11,bold=True,col=dc,al=PP_ALIGN.CENTER)
        ae=row["ae"]
        tx(s2,f"{int(ae):,}" if __import__("pandas").notna(ae) and ae>0 else "—",
           9.2,y+0.07,1.6,0.26,sz=11,col=LIGHT,al=PP_ALIGN.CENTER)
        ct=str(row["sponsor_category"]).replace("Government / ","Gov/").replace("Academic / ","Acad/")
        tx(s2,ct,10.9,y+0.07,2.1,0.26,sz=9,col=cc,al=PP_ALIGN.CENTER)
    tx(s2,"Confidential  ·  ClinicalTrials.gov  ·  2 of 3",
       0.5,7.1,12.33,0.28,sz=9,it=True,col=GRAY,al=PP_ALIGN.RIGHT)

    s3=ns()
    rc(s3,0,0,13.33,0.07,GREEN)
    tx(s3,"Strategic Recommendations",0.5,0.18,12,0.7,sz=27,bold=True)
    tx(s3,f"Five insights from the {s['condition']} landscape",
       0.5,0.9,12,0.32,sz=11,it=True,col=GRAY)
    rc(s3,0.5,1.24,12.33,0.02,BDR)
    recs=[
        (BLUE,"Pipeline Monitoring",
         f"Track {sf(s['top_sponsor'],38)} — {s['top_sponsor_trials']} studies. Monitor Phase 2→3 quarterly."),
        (RED,"Dropout Risk Mitigation",
         f"{s['dropout_rate']}% termination. Reason: {s['top_stop_reason'][:55]}."),
        (AMBER,"Geographic Expansion",
         f"{s['top_country']} dominates. Evaluate underrepresented markets."),
        (GREEN,"Phase Transition Analysis",
         f"{s['top_phase']} leads. Map the Phase 2→3 conversion funnel."),
        (PURP,"Enrollment Benchmarking",
         f"Avg {s['avg_enrollment']:,} participants. Use as baseline for design and budget."),
    ]
    for i,(color,title,body) in enumerate(recs):
        y=1.44+i*1.06
        rc(s3,0.50,y,0.05,0.88,color)
        rc(s3,0.58,y,12.2,0.88,CARD)
        tx(s3,f"0{i+1}",0.65,y+0.04,0.55,0.32,sz=13,bold=True,col=color)
        tx(s3,title,1.25,y+0.04,4.5,0.32,sz=13,bold=True,col=color)
        tx(s3,body,1.25,y+0.42,11.3,0.42,sz=11,col=LIGHT)
    tx(s3,"Confidential  ·  ClinicalTrials.gov  ·  3 of 3",
       0.5,7.1,12.33,0.28,sz=9,it=True,col=GRAY,al=PP_ALIGN.RIGHT)

    buf=io.BytesIO();prs.save(buf);buf.seek(0); return buf


# ════════════════════════════════════════════════════════
# SIDEBAR
# ════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown("""
    <div style="padding:14px 0 6px">
      <div style="font-size:1.2rem;font-weight:700;color:#e2e8f0;letter-spacing:-0.5px">
        🧬 TrialIntel
      </div>
      <div style="font-size:0.7rem;color:#1e3a5f;margin-top:2px;text-transform:uppercase;letter-spacing:1px">
        Clinical Intelligence Platform
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="sec-label">Sector</div>', unsafe_allow_html=True)
    selected_cat = st.selectbox(
        "sector", options=list(CATEGORY_CONDITIONS.keys()),
        label_visibility="collapsed"
    )

    if selected_cat != "Custom" and CATEGORY_CONDITIONS[selected_cat]:
        tags = "".join([f'<span class="tag">{s}</span>'
                        for s in CATEGORY_CONDITIONS[selected_cat]])
        st.markdown(f'<div style="margin:6px 0 10px">{tags}</div>',
                    unsafe_allow_html=True)

    st.markdown('<div class="sec-label">Condition</div>', unsafe_allow_html=True)
    condition = st.text_input("cond", value="diabetes",
                              placeholder="e.g. diabetes, alzheimer",
                              label_visibility="collapsed")

    st.markdown('<div class="sec-label">Scope</div>', unsafe_allow_html=True)
    max_results = st.select_slider(
        "scope", options=[100,200,300,500,750,1000], value=500,
        format_func=lambda x: f"{x:,} trials",
        label_visibility="collapsed"
    )

    run_ai = st.toggle("🤖 AI Analysis", value=True)
    st.markdown("")
    analyze_btn = st.button("⚡  Analyze Now", use_container_width=True)

    if "df_full" in st.session_state:
        df_full = st.session_state["df_full"]
        st.markdown('<div class="sec-label">Filters</div>', unsafe_allow_html=True)

        all_phases = sorted(df_full["phase"].dropna().unique().tolist())
        sel_phases = st.multiselect("Phase", all_phases, default=all_phases,
                                    label_visibility="collapsed",
                                    placeholder="All phases")

        all_stat = sorted(df_full["status_clean"].dropna().unique().tolist())
        sel_stat = st.multiselect("Status", all_stat, default=all_stat,
                                  label_visibility="collapsed",
                                  placeholder="All statuses")

        all_cats = sorted(df_full["sponsor_category"].dropna().unique().tolist())
        sel_cats = st.multiselect("Sponsor", all_cats, default=all_cats,
                                  label_visibility="collapsed",
                                  placeholder="All types")

        mn = int(df_full["start_year"].min()) if df_full["start_year"].notna().any() else 2000
        mx = int(df_full["start_year"].max()) if df_full["start_year"].notna().any() else 2024
        yr = st.slider("yr", mn, mx, (mn, mx),
                       label_visibility="collapsed") if mn < mx else (mn, mx)

        me = st.number_input("Min enrollment", min_value=0, value=0, step=50)

    st.markdown("---")
    st.caption("ClinicalTrials.gov · Groq Llama 3.1")


# ════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════

st.markdown("""
<div class="hero">
  <div>
    <span class="badge">● Live Data</span>
    <span class="badge">⚡ AI-Powered</span>
    <span class="badge">◈ Interactive</span>
  </div>
  <div class="hero-title">Clinical Trial Intelligence</div>
  <div class="hero-sub">
    Real-time competitive landscape analysis across phases, geographies and sponsors<br>
    Powered by ClinicalTrials.gov · Groq AI · Llama 3.1
  </div>
</div>
""", unsafe_allow_html=True)

if analyze_btn:
    if not condition.strip():
        st.error("Please enter a condition."); st.stop()

    st.markdown(f"""
    <div class="shimmer-banner">
      ⚡ Analyzing <strong>{condition.title()}</strong> —
      fetching up to <strong>{max_results:,}</strong> trials...
    </div>""", unsafe_allow_html=True)

    with st.spinner(""):
        df_raw = fetch_trials(condition.strip(), max_results)

    if df_raw.empty:
        st.error("No trials found."); st.stop()

    with st.spinner("Cleaning and enriching..."):
        df_f = clean_trials(df_raw)

    st.session_state["df_full"]   = df_f
    st.session_state["condition"] = condition.strip()
    st.session_state["summary"]   = compute_summary(df_f, condition.strip())
    st.rerun()

if "df_full" in st.session_state:
    df_full   = st.session_state["df_full"]
    condition = st.session_state["condition"]

    try:
        df = apply_filters(df_full,sel_phases,sel_stat,sel_cats,yr,me)
        if len(df) == 0: df = df_full
    except:
        df = df_full

    fs = compute_summary(df, condition)

    # KPIs
    st.markdown(f"""
    <div class="kpi-grid">
      <div class="kpi-card" style="--acc:#3b82f6">
        <span class="kpi-icon">🧪</span>
        <p class="kpi-value" style="color:#3b82f6">{fs["total_trials"]:,}</p>
        <p class="kpi-label">Trials Analyzed</p>
        <p class="kpi-delta">{len(df_full):,} total fetched</p>
      </div>
      <div class="kpi-card" style="--acc:#ef4444">
        <span class="kpi-icon">⚠️</span>
        <p class="kpi-value" style="color:#ef4444">{fs["dropout_rate"]}%</p>
        <p class="kpi-label">Dropout Rate</p>
        <p class="kpi-delta">{fs["terminated"]:,} terminated</p>
      </div>
      <div class="kpi-card" style="--acc:#10b981">
        <span class="kpi-icon">🏢</span>
        <p class="kpi-value" style="color:#10b981">{fs["unique_sponsors"]:,}</p>
        <p class="kpi-label">Sponsors</p>
        <p class="kpi-delta">{fs["industry_pct"]}% industry</p>
      </div>
      <div class="kpi-card" style="--acc:#f59e0b">
        <span class="kpi-icon">🌍</span>
        <p class="kpi-value" style="color:#f59e0b">{fs["unique_countries"]}</p>
        <p class="kpi-label">Countries</p>
        <p class="kpi-delta">Top: {fs["top_country"]}</p>
      </div>
      <div class="kpi-card" style="--acc:#8b5cf6">
        <span class="kpi-icon">👥</span>
        <p class="kpi-value" style="color:#8b5cf6">{fs["avg_enrollment"]:,}</p>
        <p class="kpi-label">Avg Enrollment</p>
        <p class="kpi-delta">Peak: {fs["peak_year"]}</p>
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(f"""
    <div class="finding-box">
      <strong>Key Finding —</strong>
      <strong>{fs["top_phase"]}</strong> trials dominate
      <strong>{condition}</strong>.
      <strong>{fs["top_sponsor"][:40]}</strong> leads with
      <strong>{fs["top_sponsor_trials"]}</strong> trials.
      Industry: <strong>{fs["industry_pct"]}%</strong>.
      Range: <strong>{fs["year_range"]}</strong>.
      Peak: <strong>{fs["peak_year"]}</strong>
      with <strong>{fs["peak_year_trials"]}</strong> starts.
    </div>
    """, unsafe_allow_html=True)

    t1,t2,t3,t4 = st.tabs([
        "📊  Dashboard","🤖  AI Analysis","📋  Data","⬇️  Downloads"
    ])

    with t1:
        cl,cr = st.columns(2, gap="medium")
        with cl:
            for fn in [chart_phase, chart_status, chart_timeline]:
                st.markdown('<div class="chart-card">', unsafe_allow_html=True)
                st.plotly_chart(fn(df), use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)
        with cr:
            for fn in [chart_sponsors, chart_map, chart_dropout]:
                st.markdown('<div class="chart-card">', unsafe_allow_html=True)
                st.plotly_chart(fn(df), use_container_width=True)
                st.markdown('</div>', unsafe_allow_html=True)
        st.markdown('<div class="chart-card">', unsafe_allow_html=True)
        st.plotly_chart(chart_enrollment(df), use_container_width=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with t2:
        if run_ai:
            with st.spinner("Generating consulting analysis..."):
                ai_raw = get_ai_analysis(fs)
            lines = []
            for line in ai_raw.split("\n"):
                if line.startswith("### "):
                    lines.append(f'<h3>{line[4:]}</h3>')
                elif line.strip():
                    lines.append(line)
            formatted = "<br>".join(lines)
            st.markdown(f'<div class="ai-box">{formatted}</div>',
                        unsafe_allow_html=True)
            st.markdown("")
            st.download_button("📄 Download as text", data=ai_raw,
                               file_name=f"analysis_{condition}.txt",
                               mime="text/plain")
        else:
            st.markdown("""
            <div class="finding-box" style="text-align:center;padding:40px">
              <span style="font-size:1.5rem">🤖</span><br><br>
              Enable <strong>AI Analysis</strong> toggle in the sidebar.
            </div>""", unsafe_allow_html=True)

    with t3:
        st.markdown(f"""
        <div style="color:#475569;font-size:0.82rem;margin-bottom:10px">
          Showing <strong style="color:#e2e8f0">{len(df):,}</strong> of
          <strong style="color:#e2e8f0">{len(df_full):,}</strong> trials
        </div>""", unsafe_allow_html=True)
        cols = ["nct_id","title","phase","status_clean","sponsor",
                "sponsor_category","primary_country","enrollment",
                "start_year","duration_months","is_terminated"]
        st.dataframe(df[[c for c in cols if c in df.columns]],
                     use_container_width=True, height=520)

    with t4:
        st.markdown("")
        ca,cb,cc = st.columns(3, gap="medium")
        with ca:
            st.markdown("""<div class="dl-card">
              <div class="dl-icon">📑</div>
              <h4>PowerPoint Deck</h4>
              <p>3-slide consulting summary with KPIs,
                 sponsor landscape and recommendations.</p>
            </div>""", unsafe_allow_html=True)
            st.markdown("")
            pptx_buf = build_pptx(df, fs)
            st.download_button("⬇️ Download PowerPoint", data=pptx_buf,
                file_name=f"trials_{condition}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True)
        with cb:
            st.markdown("""<div class="dl-card">
              <div class="dl-icon">📊</div>
              <h4>Excel Report</h4>
              <p>Multi-sheet workbook: all trials, phases,
                 top sponsors, geography.</p>
            </div>""", unsafe_allow_html=True)
            st.markdown("")
            excel_buf = build_excel(df, fs)
            st.download_button("⬇️ Download Excel", data=excel_buf,
                file_name=f"trials_{condition}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True)
        with cc:
            st.markdown("""<div class="dl-card">
              <div class="dl-icon">📄</div>
              <h4>JSON Summary</h4>
              <p>Raw intelligence metrics for further
                 analysis or API integration.</p>
            </div>""", unsafe_allow_html=True)
            st.markdown("")
            st.download_button("⬇️ Download JSON",
                data=str(fs), file_name=f"summary_{condition}.json",
                mime="application/json", use_container_width=True)
        st.markdown("---")
        st.json(fs)

else:
    st.markdown("")
    c1,c2,c3,c4 = st.columns(4, gap="medium")
    for col,(icon,title,desc) in zip(
        [c1,c2,c3,c4],[
            ("📡","Live API Data",
             "Pulls directly from ClinicalTrials.gov in real time. Up to 1,000 trials per query."),
            ("🔬","7 Analytics Charts",
             "Phase funnels, sponsor landscapes, dropout rates, maps, timelines."),
            ("🤖","AI Consulting Report",
             "Groq Llama 3.1 generates a 4-section consulting narrative automatically."),
            ("📥","Export Everything",
             "One-click PowerPoint, multi-sheet Excel, and JSON summary."),
        ]
    ):
        with col:
            st.markdown(f"""
            <div class="land-card">
              <span class="lc-icon">{icon}</span>
              <h3>{title}</h3>
              <p>{desc}</p>
            </div>""", unsafe_allow_html=True)
    st.markdown("""
    <div style="text-align:center;color:#0f2040;font-size:0.78rem;margin-top:20px;
                text-transform:uppercase;letter-spacing:1px">
      Select sector · Enter condition · Click Analyze
    </div>""", unsafe_allow_html=True)
